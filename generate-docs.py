from pptx import Presentation
from pptx.util import Inches, Pt

def create_architecture_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Azure DevOps Learning Project"
    title_slide.placeholders[1].text = "12-Day Implementation Architecture"
    
    # Slide 2: Architecture Diagram (text-based)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "GitHub → GitHub Actions CI/CD"
    
    p = tf.add_paragraph()
    p.text = "↓ Build Docker Image"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Azure Container Registry (ACR)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "↓ Deploy"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Azure Kubernetes Service (AKS)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "↓ Authenticate"
    p.level = 3
    
    p = tf.add_paragraph()
    p.text = "Microsoft Graph API (Personal Account)"
    p.level = 3
    
    # Slide 3: Technologies Used
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technologies & Services"
    
    tf = slide.placeholders[1].text_frame
    services = [
        ("Azure", ["AKS", "ACR", "Cosmos DB", "Blob Storage", "Functions", "Static Web Apps"]),
        ("CI/CD", ["GitHub Actions", "Azure DevOps Pipelines"]),
        ("Development", ["Python 3.11", "Docker", "Kubernetes", "pytest"]),
        ("APIs", ["Microsoft Graph API", "Azure SDK"])
    ]
    
    for category, items in services:
        p = tf.add_paragraph()
        p.text = category
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
    
    # Slide 4: Cost Optimization
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cost Optimization Strategy"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "$200 Azure Credits Usage:"
    
    costs = [
        ("AKS (4 days)", "$30-40"),
        ("Cosmos DB", "$25-35 (then free tier)"),
        ("VM Testing", "$30-40"),
        ("Blob Storage", "$15-20"),
        ("Total Spent", "$100-135"),
        ("Remaining", "$65-100 for experimentation")
    ]
    
    for item, cost in costs:
        p = tf.add_paragraph()
        p.text = f"{item}: {cost}"
        p.level = 1
    
    prs.save("azure-devops-architecture.pptx")
    print("✓ Presentation created: azure-devops-architecture.pptx")

if __name__ == "__main__":
    create_architecture_presentation()
